from fastapi import FastAPI, Request, HTTPException
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io

app = FastAPI()

# Helper function to convert "#008B95" to a python-pptx RGBColor object
def get_rgb(hex_code):
    hex_code = hex_code.replace("#", "") if hex_code else "000000"
    return RGBColor.from_string(hex_code)

# Helper function to apply font styling
def apply_font(shape, font_name, font_size, rgb_color, is_bold=False):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = rgb_color
            run.font.bold = is_bold

@app.post("/generate-pptx")
async def generate_pptx(request: Request):
    try:
        payload = await request.json()
        presentation_data = payload.get("presentation", {})
        slides_data = presentation_data.get("slides", [])
        theme = presentation_data.get("theme", {})
        
        # Extract Brand Theme Variables
        primary_rgb = get_rgb(theme.get("primary_color", "#008B95")) # KG Teal
        secondary_rgb = get_rgb(theme.get("secondary_color", "#1A237E")) # KG Dark Blue
        accent_rgb = get_rgb(theme.get("accent_color", "#F9A825")) # KG Orange
        charcoal_rgb = get_rgb("#212121") # Standard text
        
        font_headline = theme.get("font_headline", "Segoe UI")
        font_body = theme.get("font_body", "Open Sans")

        prs = Presentation()
        
        # Set slide dimensions to Widescreen 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_info in slides_data:
            layout_type = slide_info.get("layout", "title_slide")
            
            # --- LAYOUT 1: TITLE SLIDE ---
            if layout_type == "title_slide":
                slide = prs.slides.add_slide(prs.slide_layouts[6]) # Use blank layout for custom graphics
                
                # Draw Graphic: Dark Blue Header Banner
                header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.5))
                header_shape.fill.solid()
                header_shape.fill.fore_color.rgb = secondary_rgb
                header_shape.line.color.rgb = secondary_rgb
                
                # Draw Graphic: Orange Accent Line
                accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5), Inches(13.333), Inches(0.1))
                accent_line.fill.solid()
                accent_line.fill.fore_color.rgb = accent_rgb
                accent_line.line.color.rgb = accent_rgb
                
                # Add Title Text
                title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
                title_box.text_frame.text = slide_info.get("title", "")
                apply_font(title_box, font_headline, 44, primary_rgb, is_bold=True)
                
                # Add Subtitle Text
                sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
                sub_box.text_frame.text = slide_info.get("subtitle", "")
                apply_font(sub_box, font_body, 24, charcoal_rgb)

            # --- LAYOUT 2: SPLIT IMAGE (Left or Right) ---
            elif "image" in layout_type:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # Text area coordinates
                text_left = Inches(0.5) if "right" in layout_type else Inches(7)
                img_left = Inches(7) if "right" in layout_type else Inches(0)
                
                # Add Title
                title_box = slide.shapes.add_textbox(text_left, Inches(1), Inches(6), Inches(1))
                title_box.text_frame.text = slide_info.get("title", "")
                apply_font(title_box, font_headline, 36, secondary_rgb, is_bold=True)
                
                # Add Bullets
                body_box = slide.shapes.add_textbox(text_left, Inches(2.5), Inches(6), Inches(4))
                tf = body_box.text_frame
                bullets = slide_info.get("bullets", [])
                for i, bullet in enumerate(bullets):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = "• " + bullet
                apply_font(body_box, font_body, 20, charcoal_rgb)

                # Draw Graphic: Image Placeholder Box (with Teal border and prompt text)
                img_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, Inches(0), Inches(6.333), Inches(7.5))
                img_box.fill.solid()
                img_box.fill.fore_color.rgb = RGBColor(240, 240, 240) # Light grey
                img_box.line.color.rgb = primary_rgb
                img_box.line.width = Pt(4)
                
                # Put the image prompt text inside the box
                prompt = slide_info.get("image_prompt", "Image Placeholder")
                img_box.text_frame.text = f"[ IMAGE PLACEMENT ]\n\nPrompt: {prompt}"
                apply_font(img_box, font_body, 14, primary_rgb, is_bold=True)

            # --- LAYOUT 3: MASSIVE QUOTE (Fully colored background) ---
            elif layout_type == "massive_quote":
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # Set Background Color
                background = slide.background
                fill = background.fill
                fill.solid()
                bg_hex = slide_info.get("bg_color", "#1A237E") # Default to Dark Blue
                fill.fore_color.rgb = get_rgb(bg_hex)
                
                # Add Quote
                quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10.333), Inches(3))
                quote_box.text_frame.text = f'"{slide_info.get("quote", "")}"'
                apply_font(quote_box, font_headline, 48, get_rgb("#FFFFFF"), is_bold=True)

            # --- DEFAULT: TITLE AND BULLETS ---
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                if slide.shapes.title:
                    slide.shapes.title.text = slide_info.get("title", "")
                    apply_font(slide.shapes.title, font_headline, 36, primary_rgb, is_bold=True)
                
                if len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    for i, bullet in enumerate(slide_info.get("bullets", [])):
                        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                        p.text = bullet
                    apply_font(body, font_body, 18, charcoal_rgb)

        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        return StreamingResponse(
            ppt_stream, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=KGC_Brand_Presentation.pptx"}
        )
        
    except Exception as e:
        import traceback
        raise HTTPException(status_code=500, detail=traceback.format_exc())